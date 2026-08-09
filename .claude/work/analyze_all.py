# -*- coding: utf-8 -*-
"""分析 V1.7 全部页面的组件结构 + 表格样式详情。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation("银行流水分享_V1.7.pptx")

def shape_info(sh, indent="  "):
    x, y = sh.left / 914400, sh.top / 914400
    w, h = sh.width / 914400, sh.height / 914400
    st = str(sh.shape_type).split('(')[0].strip()
    fill = ""
    try:
        spPr = sh._element.spPr
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                fill = f"fill=#{clr.get('val')}"
    except Exception:
        pass
    # 文字
    txt = []
    try:
        for p in sh.text_frame.paragraphs:
            runs = []
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else None
                b = r.font.bold
                col = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
                runs.append(f"{r.text[:22]}({sz}pt,{'B' if b else ''},{col})")
            if runs:
                txt.append(" ".join(runs))
    except Exception:
        pass
    name = sh.name[:24]
    line = f"{indent}[{name:26s}] ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} {st} {fill}"
    if txt:
        line += "\n" + "\n".join(f"{indent}    └ {t}" for t in txt)
    return line

def table_info(sh):
    tbl = sh.table
    lines = [f"  TABLE {len(tbl.rows)}行x{len(tbl.columns)}列 pos=({sh.left/914400:.2f},{sh.top/914400:.2f}) w={sh.width/914400:.2f}in"]
    lines.append(f"    行高: {[round(r.height/914400,2) for r in tbl.rows]}")
    lines.append(f"    列宽: {[round(c.width/914400,2) for c in tbl.columns]}")
    # 第一行样式（表头）
    tc = tbl.rows[0].cells[0]
    tcPr = tc._tc.find(qn('a:tcPr'))
    fill = ""
    if tcPr is not None:
        sf = tcPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                fill = f"#{clr.get('val')}"
    # 边框
    borders = []
    if tcPr is not None:
        for tag, nm in [('a:lnB','底'), ('a:lnT','顶'), ('a:lnL','左'), ('a:lnR','右')]:
            ln = tcPr.find(qn(tag))
            if ln is not None:
                w_ = ln.get('w')
                clr = ln.find(qn('a:solidFill') + '/' + qn('a:srgbClr')) if False else None
                sf2 = ln.find(qn('a:solidFill'))
                c = None
                if sf2 is not None:
                    clr2 = sf2.find(qn('a:srgbClr'))
                    if clr2 is not None:
                        c = clr2.get('val')
                borders.append(f"{nm}{w_}/{c}")
    lines.append(f"    表头单元格: 填充{fill} 边框[{', '.join(borders)}]")
    # 字体
    rPr = None
    for p in tc.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.find(qn('a:rPr'))
            break
    if rPr is not None:
        sz = rPr.get('sz')
        b = rPr.get('b')
        latin = rPr.find(qn('a:latin'))
        ea = rPr.find(qn('a:ea'))
        col = None
        sf = rPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                col = clr.get('val')
        lines.append(f"    表头文字: {tc.text[:20]} size={int(sz)/100 if sz else None}pt bold={b} color=#{col} latin={latin.get('typeface') if latin is not None else None} ea={ea.get('typeface') if ea is not None else None}")
    # 数据行样式
    if len(tbl.rows) > 1:
        tc2 = tbl.rows[1].cells[0]
        tcPr2 = tc2._tc.find(qn('a:tcPr'))
        fill2 = ""
        if tcPr2 is not None:
            sf = tcPr2.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None:
                    fill2 = f"#{clr.get('val')}"
        bd2 = []
        if tcPr2 is not None:
            for tag, nm in [('a:lnB','底'), ('a:lnT','顶')]:
                ln = tcPr2.find(qn(tag))
                if ln is not None:
                    sf2 = ln.find(qn('a:solidFill'))
                    c = None
                    if sf2 is not None:
                        clr2 = sf2.find(qn('a:srgbClr'))
                        if clr2 is not None:
                            c = clr2.get('val')
                    bd2.append(f"{nm}{ln.get('w')}/{c}")
        lines.append(f"    数据行: 填充{fill2} 边框[{', '.join(bd2)}]")
    return "\n".join(lines)

for si, slide in enumerate(prs.slides):
    print(f"\n{'='*70}\n=== 第{si+1}页 (形状数: {len(slide.shapes)}) ===")
    for i, sh in enumerate(slide.shapes):
        if sh.shape_type == 19:
            print(table_info(sh))
        elif sh.shape_type == 6:
            print(shape_info(sh) + "\n    └ (GROUP 子元素省略)")
        else:
            print(shape_info(sh))
