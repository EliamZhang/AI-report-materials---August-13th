# -*- coding: utf-8 -*-
"""把「端到端流程与模型规划_单页汇报.pptx」第 2 页的「关键时间节点」表格
替换为水平时间轴：主线 + 6 个时间节点，上下两排事项卡片。
仅删除 tl-table 一个形状，其余元素保留。
"""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

C_BLACK    = RGBColor(0x0A, 0x0A, 0x0A)
C_DARKBLU  = RGBColor(0x1F, 0x4E, 0x79)
C_GOLD     = RGBColor(0xD6, 0xA0, 0x00)
C_BODY     = RGBColor(0x3D, 0x3C, 0x3A)
C_GRAY     = RGBColor(0x60, 0x60, 0x60)
C_LIGHTGRAY = RGBColor(0xF0, 0xF0, 0xF0)
FONT = '微软雅黑'

SRC = '端到端流程与模型规划_单页汇报.pptx'
OUT = '端到端流程与模型规划_单页汇报.pptx'

def _set_font(run, size=None, bold=None, color=None):
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    run.font.name = FONT

def _fix_ea(shape):
    for rPr in shape._element.iter(A+'rPr'):
        if rPr.find(A+'ea') is None:
            ea = etree.SubElement(rPr, A+'ea')
            ea.set('typeface', FONT)
    for pPr in shape._element.iter(A+'pPr'):
        pPr.set('latinLnBrk', '1')

def add_text(slide, x, y, w, h, text, size, bold=False, color=C_BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
             wrap=True, name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name: box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_font(r, size, bold, color)
    _fix_ea(box)
    return box

def add_rect(slide, x, y, w, h, fill=C_LIGHTGRAY, rounded=False, name=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if name: sh.name = name
    if rounded:
        try: sh.adjustments[0] = 0.08
        except Exception: pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_oval(slide, x, y, d, fill, line_color=None, name=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if name: sh.name = name
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

# ---- 数据：与第 2 页表格完全一致 ----
GROUPS = [
    ('8月9日', [
        ('943 模型平台 preview 版本完成部署，开放模型输出接口', '已完成'),
        ('下游平台开始模型调用链路开发', '进行中'),
    ]),
    ('当前', [
        ('负债 / 收入类知识库基本初始化完成，两类模型效果较好', '基本完成'),
        ('消费类知识库建设中', '进行中'),
    ]),
    ('8月21日', [
        ('完成消费类知识库初始化', '待办'),
        ('发布全分类详细模型效果评估报告', '待办'),
    ]),
    ('8月31日', [
        ('发布正式版模型', '待办'),
        ('链路调通', '待办'),
    ]),
    ('9月15日之前', [
        ('新 S-CALC 开发完成', '待办'),
        ('LMS 新页面开发完成，开 AB 前交付人审使用', '待办'),
    ]),
    ('S-CALC 开发完成后', [
        ('开通 AB 测试', '待办'),
    ]),
]

def status_color(st):
    if st in ('已完成', '基本完成'):
        return C_DARKBLU
    if st == '进行中':
        return C_GOLD
    return C_GRAY

def draw_timeline(slide):
    """无卡片时间轴：主线 + 节点 + 日期（贴节点上方）+ 事项（节点下方）。
    日期与节点紧贴，后续加时间点时复制一列即可。"""
    col_w = 1.99
    x0 = 0.56
    centers = [x0 + col_w / 2 + i * col_w for i in range(6)]
    main_y = 4.20
    date_y = 3.88
    item1_y = 4.44
    item_h = 0.30

    # 主线
    add_rect(slide, 0.72, main_y - 0.011, 11.70, 0.022, fill=C_DARKBLU, name='tl-line')

    for ci, (date, items) in enumerate(GROUPS):
        cx = centers[ci]
        x = x0 + ci * col_w
        # 主节点圆点（横跨主线）
        add_oval(slide, cx - 0.075, main_y - 0.075, 0.15, C_DARKBLU, name='tl-dot-%d' % ci)
        # 日期：紧贴节点正上方
        add_text(slide, x, date_y, col_w, 0.16, date, 8, True, C_DARKBLU, name='tl-date-%d' % ci)
        # 事项：节点下方，每项前加状态色点
        for ri, (item, st) in enumerate(items):
            y = item1_y + ri * item_h
            col = status_color(st)
            add_oval(slide, x + 0.02, y + 0.045, 0.075, col, name='tl-stdot-%d-%d' % (ci, ri))
            add_text(slide, x + 0.13, y, col_w - 0.15, item_h, item, 7, False, C_BODY,
                     line_spacing=0.95, name='tl-item-%d-%d' % (ci, ri))

def main():
    prs = Presentation(OUT)
    slide = prs.slides[1]
    # 删除原有时间轴/表格形状
    for sh in list(slide.shapes):
        if sh.name == 'tl-table' or sh.name.startswith('tl-'):
            sh._element.getparent().remove(sh._element)
    draw_timeline(slide)
    prs.save(OUT)
    print('saved', OUT)

if __name__ == '__main__':
    main()
