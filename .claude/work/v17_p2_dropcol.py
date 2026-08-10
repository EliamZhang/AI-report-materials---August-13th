# -*- coding: utf-8 -*-
"""V1.7 第2页交易明细表：删除一级分类列（中文），保留二级分类列（英文），表头改"分类"。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

tables = [sh for sh in slide.shapes if sh.shape_type == 19]
proc_shape = None
for sh in tables:
    if len(sh.table.columns) == 9:
        proc_shape = sh
        break
assert proc_shape is not None
tbl = proc_shape.table
print(f"原9列: {[round(c.width/914400,2) for c in tbl.columns]}")
print(f"表头: {[c.text for c in tbl.rows[0].cells]}")

# 校验：索引4是一级分类（中文），索引5是二级分类（英文）
hdr = [c.text for c in tbl.rows[0].cells]
assert hdr[4] == "一级分类" and hdr[5] == "二级分类", f"表头位置与预期不符: {hdr}"

# 1. 删除一级分类列（索引4）
tblGr = tbl._tbl
grid = tblGr.find(qn('a:tblGrid'))
gc_list = grid.findall(qn('a:gridCol'))
grid.remove(gc_list[4])
for tr in tblGr.findall(qn('a:tr')):
    tcs = tr.findall(qn('a:tc'))
    tr.remove(tcs[4])

# 2. 二级分类列（现索引4）表头 → "分类"
tbl.rows[0].cells[4].text = "分类"

# 3. 重新分配列宽：一级分类腾出的0.71" → 交易描述+0.45, 分类依据+0.26
#    新列宽: [0.84, 3.02, 1.34, 0.97, 1.11, 0.91, 1.00, 2.73]  合计 11.92
new_widths_emu = [int(w * 914400) for w in (0.84, 3.02, 1.34, 0.97, 1.11, 0.91, 1.00, 2.73)]
gc_list = grid.findall(qn('a:gridCol'))
for gc, w in zip(gc_list, new_widths_emu):
    gc.set('w', str(w))
for tr in tblGr.findall(qn('a:tr')):
    tcs = tr.findall(qn('a:tc'))
    for tc, w in zip(tcs, new_widths_emu):
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        tcPr.set('w', str(w))

print(f"新8列: {[round(c.width/914400,2) for c in tbl.columns]}  总宽: {sum(c.width for c in tbl.columns)/914400:.2f}in")
for ri, row in enumerate(tbl.rows):
    print(f"r{ri}: {[c.text[:16] for c in row.cells]}")

prs.save(SRC)
print("已保存")
