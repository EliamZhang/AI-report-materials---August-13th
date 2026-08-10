# -*- coding: utf-8 -*-
"""V1.7 第2页：用户汇总维度两张表加宽到全宽（左右对齐页边），竖框标签加长对齐。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

# ===== 1. 两表加宽重排 =====
# 布局：竖框标签 0.90"（含蓝条）+ 左表 x0.58 起 6.14 宽 ... 右表 x7.15 起
# 左表: x 0.58, w 5.71 (0.58-6.29)   右表: x 7.15, w 5.60 (7.15-12.75)
LEFT_TAB_X, LEFT_TAB_W = 0.58, 5.71
RIGHT_TAB_X, RIGHT_TAB_W = 7.15, 5.60
Y = 4.94
# 列宽（新）
LEFT_COLS = [2.40, 3.31]      # 指标 | 结果
RIGHT_COLS = [2.60, 3.00]     # 指标 | Fair Go Finance

tables = [sh for sh in slide.shapes if sh.shape_type == 19]
tbls = [(sh, sh.table) for sh in tables if sh.left/914400 > 2.0 and sh.width/914400 < 5.0]
assert len(tbls) == 2
tbl_left, tbl_right = tbls

def resize_table(shape, tbl, x_in, w_in, col_widths_in):
    """重设表格位置、宽度、列宽"""
    shape.left = Inches(x_in)
    shape.top = Inches(Y)
    # 1. graphicFrame ext
    xfrm = shape._element.find(qn('p:xfrm'))
    ext = xfrm.find(qn('a:ext'))
    ext.set('cx', str(int(w_in * 914400)))
    # 2. tblGrid
    tblGr = tbl._tbl
    grid = tblGr.find(qn('a:tblGrid'))
    gcs = grid.findall(qn('a:gridCol'))
    assert len(gcs) == 2
    for gc, w in zip(gcs, col_widths_in):
        gc.set('w', str(int(w * 914400)))
    # 3. 每行 tcPr w
    for tr in tblGr.findall(qn('a:tr')):
        tcs = tr.findall(qn('a:tc'))
        for tc, w in zip(tcs, col_widths_in):
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            tcPr.set('w', str(int(w * 914400)))

resize_table(*tbl_left, LEFT_TAB_X, LEFT_TAB_W, LEFT_COLS)
resize_table(*tbl_right, RIGHT_TAB_X, RIGHT_TAB_W, RIGHT_COLS)
print(f"左表 → ({LEFT_TAB_X},{Y}) {LEFT_TAB_W}in 列宽{LEFT_COLS}")
print(f"右表 → ({RIGHT_TAB_X},{Y}) {RIGHT_TAB_W}in 列宽{RIGHT_COLS}")

# ===== 2. 竖框标签加长到与表格同高（2.51in）=====
# 原标签：左(0.90,5.22) 右(7.07,5.24)，高1.47
# 新位置：左(0.90,4.94) 右(7.07,4.94)，高2.51
for sh in slide.shapes:
    if sh.name.startswith("p2-label"):
        x, y = sh.left/914400, sh.top/914400
        if abs(x - 0.90) < 0.05 or abs(x - 7.07) < 0.05:
            sh.top = Inches(4.94)
            sh.height = Inches(2.51)
            # 蓝条也加长
            print(f"{sh.name} → ({sh.left/914400:.2f},{sh.top/914400:.2f}) 高{sh.height/914400:.2f}")

prs.save(SRC)
print("已保存")
