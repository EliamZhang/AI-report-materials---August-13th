# -*- coding: utf-8 -*-
"""读取V1.7第2页两张表的交易描述列完整内容。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[1]

for i, sh in enumerate(slide.shapes):
    if sh.shape_type != 19:
        continue
    tbl = sh.table
    print(f"=== Table (形状{i}) {len(tbl.rows)}行x{len(tbl.columns)}列 ===")
    for ri, row in enumerate(tbl.rows):
        cells = [c.text for c in row.cells]
        print(f"r{ri:2d}: {cells}")
