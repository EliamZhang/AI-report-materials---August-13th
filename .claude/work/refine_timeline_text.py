# -*- coding: utf-8 -*-
"""梳理单页汇报时间轴下方事项文字：统一术语、规范中英混排、动词一致。
只改文字（保留段落结构、字号、颜色、位置），不动排版。
"""
from pptx import Presentation

SRC = '端到端流程与模型规划_单页汇报.pptx'

# (文本框名, 段落索引, 新文本)；tl-item-3-0 重名，用 x 坐标过滤
M = [
    ('tl-item-0-0', None, 1, '发布 Liability 识别报告'),
    ('tl-item-0-0', None, 2, '下游平台启动模型调用链路开发'),
    ('tl-item-1-0', None, 0, '负债/收入类知识库初始化基本完成'),
    ('tl-item-1-0', None, 2, '新 S-CALC 开发中'),
    ('tl-item-1-0', None, 3, '启动新旧 IDP 链路 AB 测试'),
    ('tl-item-2-0', None, 1, '发布全分类模型效果详细评估报告'),
    ('tl-item-3-0', 7.0, 1, '新分类数据落库'),
    ('tl-item-3-0', 7.0, 3, '链路调通后空跑两周'),
    ('tl-item-5-0', None, 0, '启动新旧 IDP 链路 AB 测试'),
    ('tl-item-5-0', None, 3, '完成新旧 IDP 链路 AB 测试'),
]

p = Presentation(SRC)
s = p.slides[0]

def find_shape(name, x_lo):
    for sh in s.shapes:
        if sh.name != name or not sh.has_text_frame:
            continue
        l = sh.left / 914400
        if x_lo is not None and not (x_lo - 0.3 <= l <= x_lo + 0.3):
            continue
        return sh
    raise KeyError('not found: %s @%s' % (name, x_lo))

def set_para(shape, idx, new_text):
    para = shape.text_frame.paragraphs[idx]
    runs = para.runs
    assert runs, 'empty para %d in %s' % (idx, shape.name)
    first = runs[0]
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    first.text = new_text

for name, x_lo, idx, new in M:
    set_para(find_shape(name, x_lo), idx, new)

p.save(SRC)
print('done')
