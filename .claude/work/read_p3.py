# -*- coding: utf-8 -*-
"""读取V1.7第3页右侧区域：监控/定位/输出文字块的结构。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[2]

print("=== 第3页右侧区域形状 (x>=6.0) ===")
for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    if x >= 5.8 or w > 8:
        print(f"\n[{i}] {sh.name} ({sh.shape_type}) pos=({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} 底={y+h:.2f}")
        if sh.has_text_frame:
            tf = sh.text_frame
            for pi, p in enumerate(tf.paragraphs):
                runs = []
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else None
                    b = r.font.bold
                    col = None
                    try:
                        col = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
                    except Exception:
                        col = "theme"
                    runs.append(f"[{r.text}({sz}pt,{'B' if b else ''},{col})]")
                print(f"  P{pi}: {' '.join(runs)}")
