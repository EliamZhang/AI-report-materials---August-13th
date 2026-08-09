# -*- coding: utf-8 -*-
"""分析 V1.7 第1页完整结构：形状、位置、字号、颜色、间距。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[0]

def fmt_runs(shape):
    """提取文本框内的文字、字号、粗体、颜色"""
    out = []
    try:
        tf = shape.text_frame
        for p in tf.paragraphs:
            runs = []
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else None
                b = r.font.bold
                col = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
                runs.append(f"{r.text[:18]}({sz}pt,{'B' if b else ''},{col})")
            out.append(" | ".join(runs))
    except Exception:
        pass
    return out

print(f"=== V1.7 第1页 形状数: {len(slide.shapes)} ===")
for i, sh in enumerate(slide.shapes):
    x, y = sh.left / 914400, sh.top / 914400
    w, h = sh.width / 914400, sh.height / 914400
    st = str(sh.shape_type).split('(')[0].strip()
    name = sh.name
    fill = ""
    try:
        spPr = sh._element.spPr
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                fill = f"fill={clr.get('val')}"
                a = clr.find(qn('a:alpha'))
                if a is not None:
                    fill += f" alpha={int(a.get('val'))/1000:.2f}"
    except Exception:
        pass
    print(f"[{i:2d}] {name[:26]:28s} ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} {st} {fill}")
    txt = fmt_runs(sh)
    if txt:
        for t in txt:
            print(f"      └ {t}")
    # 表格
    if sh.shape_type == 19:
        tbl = sh.table
        print(f"      └ TABLE {len(tbl.rows)}行 x {len(tbl.columns)}列, 行高={[round(r.height/914400,2) for r in tbl.rows]}")
        for ri, row in enumerate(tbl.rows[:3]):
            cells = []
            for c in row.cells:
                t = c.text.replace("\n", " ")[:14]
                cells.append(t)
            print(f"        r{ri}: {cells}")
