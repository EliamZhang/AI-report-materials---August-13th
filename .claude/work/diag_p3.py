# -*- coding: utf-8 -*-
"""诊断：扫描文字块区域的逐行深色像素分布，找出行结构。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
import numpy as np
from PIL import Image

img = np.array(Image.open(r".claude\work\p3_check\slide3.png").convert("RGB")).astype(int)
H, W, _ = img.shape
SX, SY = W/13.333, H/7.5

def near(c, t, tol=80):
    return all(abs(int(c[i])-t[i])<=tol for i in range(3))

x0, x1 = int(6.9*SX), int(12.96*SX)
BODY = (61, 60, 58)

# 逐行扫描 y 2.0-4.6in
lines = []
for yy in range(int(2.0*SY), int(4.6*SY)):
    row = img[yy, x0:x1]
    dark = sum(near(p, BODY) for p in row)
    if dark > 3:
        lines.append((yy, dark))

# 合并连续行成文本带
bands = []
for yy, d in lines:
    if bands and yy - bands[-1][1] <= 2:
        bands[-1][1] = yy
        bands[-1][2] = max(bands[-1][2], d)
    else:
        bands.append([yy, yy, d])

print(f"检测到 {len(bands)} 个文本带:")
for i, (a, b, d) in enumerate(bands):
    print(f"  带{i}: y {a/SY:.2f}-{b/SY:.2f}in (高{(b-a)/SY:.2f}in) 深色峰值={d}")

# XML: 检查文本框 lineSpacing 和 autoSize
from pptx import Presentation
from pptx.oxml.ns import qn
prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[2]
for sh in slide.shapes:
    if abs(sh.left/914400 - 6.90) < 0.05 and abs(sh.top/914400 - 2.16) < 0.05:
        bodyPr = sh._element.txBody.find(qn('a:bodyPr'))
        print("\nbodyPr:", bodyPr.attrib if bodyPr is not None else None)
        for pi, p in enumerate(sh.text_frame.paragraphs):
            pPr = p._p.find(qn('a:pPr'))
            lnSp = pPr.find(qn('a:lnSpc')) if pPr is not None else None
            spA = pPr.find(qn('a:spcAft')) if pPr is not None else None
            ln = None
            if lnSp is not None:
                spc = lnSp.find(qn('a:spcPct'))
                ln = spc.get('val') if spc is not None else None
            aft = None
            if spA is not None:
                spc = spA.find(qn('a:spcPts'))
                aft = spc.get('val') if spc is not None else None
            print(f"P{pi}: lineSpacing={ln} spaceAfter={aft} text={p.runs[0].text[:15] if p.runs else ''}")
        break
