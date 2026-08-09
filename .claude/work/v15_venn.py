# -*- coding: utf-8 -*-
"""V1.5 第4页：把左侧条形图替换为真正的韦恩图（两圆相交+数字标签）。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.5.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[3]

BLUE = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xD6, 0xA0, 0x00)
TEXT_DARK = RGBColor(0x0A, 0x0A, 0x0A)
TEXT_BODY = RGBColor(0x3D, 0x3C, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def delete_shape(sh):
    spTree = slide.shapes._spTree
    spTree.remove(sh._element)

# ============ 1. 删除旧条形图 ============
chart = None
for sh in slide.shapes:
    if sh.shape_type == 3:
        chart = sh
        break
if chart is not None:
    delete_shape(chart)
    print("旧条形图已删除")

# ============ 2. 绘制韦恩图（两圆相交）============
def add_oval_alpha(x, y, d, rgb, alpha_pct, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    # 设置透明度
    spPr = sh._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    clr = sf.find(qn('a:srgbClr'))
    alpha = etree.SubElement(clr, qn('a:alpha'))
    alpha.set('val', str(alpha_pct * 1000))
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1.2)
    sh.shadow.inherit = False
    return sh

def add_label(x, y, w, h, text, size, bold, color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
        p.alignment = align
        p.space_after = Pt(0)
        p.line_spacing = 1.0
    return tb

# 几何参数（英寸）
D = 2.35                 # 圆直径
Y0 = 1.55                # 圆顶
LEFT_X = 0.55            # 左圆
RIGHT_X = 1.85           # 右圆（圆心距 1.3，重叠 1.05）
CY = Y0 + D / 2          # 圆心 y = 2.725

# 左圆（Finv 蓝色）
add_oval_alpha(LEFT_X, Y0, D, BLUE, 28, line_color=BLUE)
# 右圆（illion 金色）
add_oval_alpha(RIGHT_X, Y0, D, GOLD, 28, line_color=GOLD)

# 圆内数字标签
# 左独有 396：左圆左侧
add_label(LEFT_X - 0.15, CY - 0.15, 1.0, 0.30, "396", 11, True, BLUE)
# 右独有 79：右圆右侧
add_label(RIGHT_X + D - 0.85, CY - 0.15, 1.0, 0.30, "79", 11, True, GOLD)
# 交集 5,284：中心
add_label(LEFT_X + 0.95, CY - 0.22, 0.75, 0.30, "5,284", 12, True, TEXT_DARK)
# 交集副标签：一致率
add_label(LEFT_X + 0.78, CY + 0.10, 1.05, 0.22, "一致率 92.0%", 7.5, False, TEXT_BODY)

# 圆外标签（顶部）
add_label(LEFT_X - 0.25, Y0 - 0.55, D + 0.5, 0.30, "Finv 识别 5,680", 9, True, BLUE)
add_label(LEFT_X - 0.25, Y0 - 0.28, D + 0.5, 0.24, "覆盖率 12.8%", 8, False, TEXT_BODY)
add_label(RIGHT_X - 0.25, Y0 - 0.55, D + 0.5, 0.30, "illion 识别 5,363", 9, True, GOLD)
add_label(RIGHT_X - 0.25, Y0 - 0.28, D + 0.5, 0.24, "覆盖率 12.1%", 8, False, TEXT_BODY)

print("韦恩图绘制完成")

prs.save(SRC)
print("已保存")
