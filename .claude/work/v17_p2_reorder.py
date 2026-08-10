# -*- coding: utf-8 -*-
"""V1.7 第2页交易明细表：列顺序调整为 日期→交易金额→交易描述→其余，列宽同步重排。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

tables = [sh for sh in slide.shapes if sh.shape_type == 19]
proc_shape = None
for sh in tables:
    if len(sh.table.columns) == 8:
        proc_shape = sh
        break
assert proc_shape is not None
tbl = proc_shape.table

hdr = [c.text for c in tbl.rows[0].cells]
print(f"原列序: {hdr}")
assert hdr[0] == "日期" and hdr[1] == "交易描述" and hdr[6] == "交易金额"

# 新列序：日期, 交易金额, 交易描述, 标准交易对手, 主体类型, 分类, 交易性质, 分类依据
order = [0, 6, 1, 2, 3, 4, 5, 7]

tblGr = tbl._tbl

# 1. 重排 gridCol
grid = tblGr.find(qn('a:tblGrid'))
gc_list = grid.findall(qn('a:gridCol'))
widths = [gc.get('w') for gc in gc_list]
new_widths = [widths[i] for i in order]
for gc, w in zip(gc_list, new_widths):
    gc.set('w', w)

# 2. 重排每行 tc + 同步 tcPr w
for tr in tblGr.findall(qn('a:tr')):
    tcs = tr.findall(qn('a:tc'))
    new_tcs = [tcs[i] for i in order]
    for tc in new_tcs:
        tr.remove(tc)
    for tc, w in zip(new_tcs, new_widths):
        tr.append(tc)
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        tcPr.set('w', w)

print(f"新列序: {[c.text for c in tbl.rows[0].cells]}")
print(f"新列宽: {[round(c.width/914400,2) for c in tbl.columns]}  总宽: {sum(c.width for c in tbl.columns)/914400:.2f}in")
for ri in range(min(2, len(tbl.rows))):
    row = tbl.rows[ri]
    print(f"r{ri}: {[c.text[:16] for c in row.cells]}")

prs.save(SRC)
print("已保存")
