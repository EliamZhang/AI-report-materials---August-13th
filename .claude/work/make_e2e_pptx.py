# -*- coding: utf-8 -*-
"""生成「端到端流程与模型规划」两版单页 PPT（V2.0 风格）：
布局 A：上流程 + 下闭环；布局 B：左流程 + 右闭环（双栏）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

C_BLACK   = RGBColor(0x0A, 0x0A, 0x0A)
C_DARKBLU = RGBColor(0x1F, 0x4E, 0x79)
C_GOLD    = RGBColor(0xD6, 0xA0, 0x00)
C_BODY    = RGBColor(0x3D, 0x3C, 0x3A)
C_GRAY    = RGBColor(0x60, 0x60, 0x60)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHTGOLD = RGBColor(0xFA, 0xF3, 0xE0)
C_LIGHTGRAY = RGBColor(0xF0, 0xF0, 0xF0)
FONT = '微软雅黑'

def _set_font(run, size=None, bold=None, color=None):
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    run.font.name = FONT

def _fix_ea(shape):
    """补 ea 东亚字体 + latinLnBrk=1，避免中文渲染溢出（规范 34.1）"""
    for rPr in shape._element.iter(A+'rPr'):
        if rPr.find(A+'ea') is None:
            ea = etree.SubElement(rPr, A+'ea')
            ea.set('typeface', FONT)
    for pPr in shape._element.iter(A+'pPr'):
        pPr.set('latinLnBrk', '1')

def add_text(slide, x, y, w, h, text, size, bold=False, color=C_BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
             wrap=True, name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name: box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_font(r, size, bold, color)
    _fix_ea(box)
    return box

def add_runs(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None, wrap=True, name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name: box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    for t, s, b, c in runs:
        r = p.add_run(); r.text = t
        _set_font(r, s, b, c)
    _fix_ea(box)
    return box

def add_bar(slide, x, y, w=0.08, h=0.25, color=C_GOLD, name=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if name: sh.name = name
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_rect(slide, x, y, w, h, fill=C_LIGHTGRAY, rounded=False, name=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if name: sh.name = name
    if rounded:
        try: sh.adjustments[0] = 0.05
        except Exception: pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_oval_num(slide, x, y, d, num, size=9, fill=C_DARKBLU, name=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if name: sh.name = name
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    _set_font(r, size, True, C_WHITE)
    _fix_ea(sh)
    return sh

def add_card(slide, x, y, w, h, title_runs, items, item_h=0.30, name_size=8,
             det_size=7, name=None):
    """浅灰圆角卡：深蓝顶条 + 标题 + 条目（编号+名称 / 细节两行）"""
    add_rect(slide, x, y, w, h, rounded=True, name=(name or 'card')+'-bg')
    add_rect(slide, x, y, w, 0.06, fill=C_DARKBLU, name=(name or 'card')+'-topbar')
    add_runs(slide, x+0.15, y+0.10, w-0.30, 0.22, title_runs, name=(name or 'card')+'-title')
    iy = y + 0.36
    for num, nm, det in items:
        add_runs(slide, x+0.15, iy, w-0.30, 0.16,
                 [(num + ' ', name_size, True, C_DARKBLU), (nm, name_size, True, C_BLACK)],
                 name=(name or 'card')+'-it')
        add_text(slide, x+0.34, iy+0.17, w-0.46, 0.14, det, det_size, color=C_BODY,
                 name=(name or 'card')+'-det')
        iy += item_h

def add_plain_table(slide, x, y, col_widths, row_heights, data, header_font=8,
                    body_font=8, header_color=C_BLACK, body_color=C_BODY,
                    status_col=None, name=None):
    """极简横线表：无左右边框，表头底 1.5pt 黑线，数据行底 0.75pt 灰线（规范 8.1/34.8）"""
    nrows, ncols = len(data), len(col_widths)
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                Inches(sum(col_widths)), Inches(sum(row_heights)))
    tbl = gf.table
    if name: gf.name = name
    tblPr = tbl._tbl.find(A+'tblPr')
    if tblPr is not None:
        for el in list(tblPr):
            if el.tag == A+'tableStyleId':
                tblPr.remove(el)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    for ri, rh in enumerate(row_heights):
        tbl.rows[ri].height = Inches(rh)
    for ri in range(nrows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            tcPr = cell._tc.find(A+'tcPr')
            if tcPr is None:
                tcPr = etree.SubElement(cell._tc, A+'tcPr')
            for ln in ['lnL','lnR','lnT','lnB']:
                e = tcPr.find(A+ln)
                if e is not None: tcPr.remove(e)
            tcPr.set('marT','0'); tcPr.set('marB','0')
            tcPr.set('anchor','ctr')
            def _no_fill(tag):
                ln = etree.SubElement(tcPr, A+tag)
                etree.SubElement(ln, A+'noFill')
            def _line(tag, w, color):
                ln = etree.SubElement(tcPr, A+tag)
                ln.set('w', str(w))
                sf = etree.SubElement(ln, A+'solidFill')
                etree.SubElement(sf, A+'srgbClr').set('val', color)
            _no_fill('lnL'); _no_fill('lnR')
            if ri == 0:
                _no_fill('lnT')
                _line('lnB', 19050, '0A0A0A')
            else:
                _no_fill('lnT')
                _line('lnB', 9525, 'D6D6D6')
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.03)
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 0.9
            r = p.add_run(); r.text = data[ri][ci] if data[ri][ci] else ' '
            if ri == 0:
                _set_font(r, header_font, True, header_color)
            elif status_col is not None and ci == status_col:
                txt = data[ri][ci]
                if txt in ('已完成', '基本完成'):
                    _set_font(r, body_font, True, C_DARKBLU)
                elif txt == '进行中':
                    _set_font(r, body_font, False, C_GRAY)
                else:
                    _set_font(r, body_font, False, body_color)
            else:
                _set_font(r, body_font, False, body_color)
            _fix_ea(gf)
    return gf

def add_conclusion(slide, x, y, w, text, size=8.5, h=0.34, name=None):
    add_rect(slide, x, y, w, h, fill=C_LIGHTGOLD, name=(name or 'concl')+'-bg')
    add_rect(slide, x, y, 0.07, h, fill=C_GOLD, name=(name or 'concl')+'-bar')
    add_text(slide, x+0.20, y, w-0.35, h, text, size, True, C_BLACK,
             anchor=MSO_ANCHOR.MIDDLE, name=(name or 'concl')+'-text')

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def header(slide, title, subtitle):
    add_text(slide, 0.66, 0.23, 9.72, 0.39, title, 20, True, C_BLACK, name='title')
    add_text(slide, 0.56, 0.64, 11.92, 0.48, subtitle, 18, False, C_GRAY, name='subtitle')

# ============================================================
# 内容
# ============================================================
TITLE = 'bs_cat 已可线上调用：端到端链路 + 1 人 2 Agent 运维闭环'
SUBTITLE = '从申请到输出全链路自主可控；8 月底调通、9 月中旬新 S-CALC 与新页面上线后开 AB'

FLOW = [
    ('阶段 1｜申请与调度', [
        ('1', 'IDA 发起申请', '发起贷款申请与流水处理请求'),
        ('2', '进入消息队列', 'AWS SQS 接收并排队任务'),
        ('3', 'IDP 核心调度', '拉取申请、预处理数据、编排流程'),
    ]),
    ('阶段 2｜分类与计算', [
        ('4', '调用流水分类模型', '调用 BS-CAT 分类能力'),
        ('5', '交易分类与识别', '输出交易分类、交易对手、分析标签'),
        ('6', '服务能力计算', 'S-CALC 执行收入/负债/支出等计算'),
        ('7', 'GenAI 辅助推理', '处理复杂规则与补充判断'),
        ('8', '结果落库', '保存分类结果、HTM 流水、测算结果'),
    ]),
    ('阶段 3｜决策与输出', [
        ('9', '风险模型与决策', '结合画像与历史表现，输出评分/建议'),
        ('10', 'LMS 输出应用', '返回推荐额度、质量标签、可视化结果'),
    ]),
]
SUPPORT = '支撑能力：A. 分类知识库（规则、关键词、第三方数据）→ 支撑交易分类   ｜   B. 分类模型监控与运营（效果监控、异常跟踪、人工维护）→ 支撑风险决策'
TRADE_A = '核心取舍：分类模型直接输出“可用于信审”的核心结果；IDA 只保留必要的调度、汇总与基础计算，降低下游复杂度'
TRADE_B = '核心取舍：模型直接输出“可用于信审”的结果；IDA 只保留调度、汇总与基础计算'

TIMELINE = [
    ('时间节点', '事项', '状态'),
    ('8月9日', '943 模型平台 preview 版本完成部署，开放模型输出接口', '已完成'),
    ('8月9日', '下游平台开始模型调用链路开发', '进行中'),
    ('当前', '负债 / 收入类知识库基本初始化完成，两类模型效果较好', '基本完成'),
    ('当前', '消费类知识库建设中', '进行中'),
    ('8月21日', '完成消费类知识库初始化', '待办'),
    ('8月21日', '发布全分类详细模型效果评估报告', '待办'),
    ('8月31日', '发布正式版模型', '待办'),
    ('8月31日', '链路调通', '待办'),
    ('9月15日之前', '新 S-CALC 开发完成', '待办'),
    ('9月15日之前', 'LMS 新页面开发完成，开 AB 前交付人审使用', '待办'),
    ('S-CALC 开发完成后', '开通 AB 测试', '待办'),
]
DIVISION = [
    ('调用方', '负责人', '职责'),
    ('IDA', 'Tien / Niral', '链路调用、调度逻辑、Serviceability 计算与 IDA/IDP 代码落地'),
    ('风险 / 模型', 'Risk & Model', '使用分类、负债、收入、支出与 Serviceability 指标'),
    ('人审', 'SAV', 'Case 复核、异常解释、质检样本与业务反馈闭环'),
]
AGENT_QC = [
    ('①', '效果监控', '分类分布与命中率、数据/分类漂移；对手方识别异常、关键字段缺失'),
    ('②', '问题诊断', '定位异常交易与问题引擎，识别缺失商户/关键词/规则；自动生成优先级 Todo 清单'),
]
AGENT_KB = [
    ('①', '知识补全', '检索未知商户/对手方，补充企业、品牌、关键词与分类信息，识别新增业务实体'),
    ('②', '知识校验', '交叉验证分类与商户信息，识别冲突、过期及低置信知识，输出待人工确认项'),
    ('③', '知识库更新', '自动生成新增/修改建议，人工审核后进入正式知识库，持续沉淀规则与案例'),
]
CLOSED = '形成“监控 → 定位 → 修复 → 沉淀”持续运维闭环，线上问题持续沉淀为模型能力'

# ============================================================
# 布局 A：上流程 + 下闭环
# ============================================================
def build_A(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = new_slide(prs)
    header(s, TITLE, SUBTITLE)

    # 章节 1
    add_bar(s, 0.59, 1.28)
    add_text(s, 0.78, 1.26, 10.0, 0.25, '1. 端到端链路：从申请发起到信审输出', 14, True, C_BLACK, name='s1-title')
    card_w, gap, y0, card_h = 3.90, 0.25, 1.62, 1.46
    xs = [0.56, 0.56 + card_w + gap, 0.56 + 2*(card_w + gap)]
    for (tname, steps), x in zip(FLOW, xs):
        add_rect(s, x, y0, card_w, card_h, rounded=True, name='stage-bg-'+tname[0])
        add_rect(s, x, y0, card_w, 0.06, fill=C_DARKBLU, name='stage-topbar-'+tname[0])
        add_text(s, x+0.15, y0+0.11, card_w-0.30, 0.22, tname, 10.5, True, C_DARKBLU, name='stage-title-'+tname[0])
        iy = y0 + 0.40
        for num, nm, det in steps:
            add_oval_num(s, x+0.14, iy+0.02, 0.22, num, size=8.5, name='num-'+num)
            add_runs(s, x+0.45, iy, card_w-0.60, 0.18,
                     [(nm + '｜', 8.5, True, C_BLACK), (det, 7.5, False, C_GRAY)], name='step-'+num)
            iy += 0.205
    add_text(s, 4.50, y0+0.50, 0.20, 0.3, '→', 14, True, C_DARKBLU, align=PP_ALIGN.CENTER, name='arrow1')
    add_text(s, 8.65, y0+0.50, 0.20, 0.3, '→', 14, True, C_DARKBLU, align=PP_ALIGN.CENTER, name='arrow2')
    add_text(s, 0.56, 3.20, 12.2, 0.20, SUPPORT, 8, False, C_GRAY, name='support')
    add_conclusion(s, 0.56, 3.48, 12.2, TRADE_A, name='trade')

    # 章节 2
    add_bar(s, 0.59, 4.02)
    add_text(s, 0.78, 4.00, 10.0, 0.25, '2. 部署与运维闭环：1 人 + 2 个 AI Agent', 14, True, C_BLACK, name='s2-title')
    add_text(s, 0.56, 4.38, 5.5, 0.22, '关键时间节点', 12, True, C_BLACK, name='tl-title')
    add_plain_table(s, 0.56, 4.70, [1.25, 3.60, 0.65], [0.17]+[0.15]*11,
                    TIMELINE, header_font=8, body_font=8, status_col=2, name='tl-table')
    add_text(s, 0.56, 6.66, 5.50, 0.42,
             '下游分工：IDA（Tien / Niral）链路调用、调度与代码落地 ｜ 风险 / 模型使用分类与 Serviceability 指标 ｜ 人审（SAV）Case 复核与反馈闭环',
             8, False, C_GRAY, name='division', line_spacing=1.1)
    add_text(s, 6.28, 4.38, 6.48, 0.22, '1 人 + 2 个 AI Agent：监控 → 定位 → 修复 → 沉淀', 12, True, C_BLACK, name='agent-title')
    add_card(s, 6.28, 4.70, 6.48, 0.92,
             [('AI Agent 1｜', 10, True, C_DARKBLU), ('质检 Agent：监控异常、定位问题', 10, True, C_BLACK)],
             AGENT_QC, item_h=0.30, name='agent-qc')
    add_card(s, 6.28, 5.70, 6.48, 1.22,
             [('AI Agent 2｜', 10, True, C_DARKBLU), ('知识库 Agent：主动补充、持续迭代', 10, True, C_BLACK)],
             AGENT_KB, item_h=0.30, name='agent-kb')
    add_conclusion(s, 6.28, 7.00, 6.48, CLOSED, size=8, h=0.30, name='closed')
    prs.save(path)
    print('saved', path)

# ============================================================
# 布局 B：左流程 + 右闭环（双栏）
# ============================================================
def build_B(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = new_slide(prs)
    header(s, TITLE, SUBTITLE)

    # 左栏：端到端链路
    add_bar(s, 0.56, 1.28)
    add_text(s, 0.75, 1.26, 4.8, 0.25, '1. 端到端链路', 14, True, C_BLACK, name='s1-title')
    y = 1.64
    for num, nm, det in [st for _, steps in FLOW for st in steps]:
        add_oval_num(s, 0.56, y+0.02, 0.24, num, size=9, name='num-'+num)
        add_runs(s, 0.92, y, 4.45, 0.20,
                 [(nm, 9.5, True, C_BLACK), ('  ｜  ', 7.5, False, C_GRAY), (det, 7.5, False, C_GRAY)],
                 name='step-'+num)
        y += 0.40
    add_text(s, 0.56, y+0.06, 4.84, 0.20, '支撑能力：A. 分类知识库（规则、关键词、第三方数据）→ 支撑交易分类', 8, False, C_GRAY, name='support1')
    add_text(s, 0.56, y+0.30, 4.84, 0.20, 'B. 分类模型监控与运营（效果监控、异常跟踪、人工维护）→ 支撑风险决策', 8, False, C_GRAY, name='support2')
    add_conclusion(s, 0.56, y+0.62, 4.84, TRADE_B, size=7.5, h=0.42, name='trade')

    # 右栏：部署与运维闭环
    add_bar(s, 6.28, 1.28)
    add_text(s, 6.47, 1.26, 6.4, 0.25, '2. 部署与运维闭环', 14, True, C_BLACK, name='s2-title')
    add_text(s, 6.28, 1.60, 6.4, 0.20, '关键时间节点', 11, True, C_BLACK, name='tl-title')
    add_plain_table(s, 6.28, 1.90, [1.30, 4.40, 0.70], [0.16]+[0.145]*11,
                    TIMELINE, header_font=7.5, body_font=7.5, status_col=2, name='tl-table')
    add_text(s, 6.28, 3.80, 6.4, 0.20, '下游分工', 11, True, C_BLACK, name='div-title')
    add_plain_table(s, 6.28, 4.06, [1.00, 1.15, 4.25], [0.16]+[0.19]*3,
                    DIVISION, header_font=7.5, body_font=7.5, name='div-table')
    add_text(s, 6.28, 4.96, 6.4, 0.20, '1 人 + 2 个 AI Agent：监控 → 定位 → 修复 → 沉淀', 11, True, C_BLACK, name='agent-title')
    add_card(s, 6.28, 5.24, 6.48, 0.92,
             [('AI Agent 1｜', 9.5, True, C_DARKBLU), ('质检 Agent：监控异常、定位问题', 9.5, True, C_BLACK)],
             AGENT_QC, item_h=0.30, name='agent-qc')
    add_card(s, 6.28, 6.24, 6.48, 1.22,
             [('AI Agent 2｜', 9.5, True, C_DARKBLU), ('知识库 Agent：主动补充、持续迭代', 9.5, True, C_BLACK)],
             AGENT_KB, item_h=0.30, name='agent-kb')
    prs.save(path)
    print('saved', path)

if __name__ == '__main__':
    build_A('端到端流程与模型规划_布局A_上流程下闭环.pptx')
    build_B('端到端流程与模型规划_布局B_左流程右闭环.pptx')
