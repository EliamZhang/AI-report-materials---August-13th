# -*- coding: utf-8 -*-
"""读取第2页三个小节标题的精确位置与样式。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[1]

for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    txt = sh.text_frame.text[:40] if sh.has_text_frame else ""
    if txt.startswith(("1. ", "2. ", "3. ")) or sh.name in ("s1-title", "s2-title", "s3-title"):
        print(f"[{i:2d}] {sh.name[:20]:22s} ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} 底={y+h:.2f}  {txt}")
