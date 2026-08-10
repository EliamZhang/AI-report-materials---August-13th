# -*- coding: utf-8 -*-
"""V1.7 第3页：文字块上移+加高，容纳拆行后的8行内容，避免压到下方标题。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[2]

target = None
for sh in slide.shapes:
    if abs(sh.left/914400 - 6.90) < 0.05 and abs(sh.top/914400 - 2.18) < 0.05 and sh.has_text_frame:
        target = sh
        break
assert target is not None

old_top = target.top / 914400
old_h = target.height / 914400
# 上移到 2.16，加高到 1.96（底 4.12，标题 4.13 上方）
target.top = Inches(2.16)
target.height = Inches(1.96)
print(f"文字块: y {old_top:.2f}→{target.top/914400:.2f}, 高 {old_h:.2f}→{target.height/914400:.2f}, 底 {old_top+old_h:.2f}→{target.top/914400+target.height/914400:.2f}")

prs.save(SRC)
print("已保存")
