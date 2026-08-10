# -*- coding: utf-8 -*-
"""V1.7 第2页：用户汇总维度两个竖框标签美化为 浅蓝底+左侧蓝条+竖排文字。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

BLUE = RGBColor(0x1F, 0x4E, 0x79)
BLUE_LIGHT = RGBColor(0xE3, 0xEB, 0xF4)
TEXT_DARK = RGBColor(0x0A, 0x0A, 0x0A)

# 两个竖框的位置（矩形10、11）
targets = [
    (0.90, 5.22, "用户维度"),      # 左表格
    (7.07, 5.24, "单笔借款维度"),  # 右表格
]

for x, y, text in targets:
    # 1. 找到原矩形并删除
    found = None
    for sh in slide.shapes:
        if sh.shape_type == 1 and abs(sh.left/914400 - x) < 0.05 and abs(sh.top/914400 - y) < 0.05:
            found = sh
            break
    if found is not None:
        slide.shapes._spTree.remove(found._element)

    # 2. 浅蓝底矩形
    bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.85), Inches(1.47))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_LIGHT
    bg.line.fill.background()
    bg.shadow.inherit = False
    bg.name = "p2-label-bg"

    # 3. 左侧蓝条
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(1.47))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.name = "p2-label-bar"

    # 4. 竖排文字
    tb = slide.shapes.add_textbox(Inches(x + 0.06), Inches(y), Inches(0.79), Inches(1.47))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 竖排（eaVert）：中文从上到下排列
    bodyPr = tb._element.txBody.find(qn('a:bodyPr'))
    bodyPr.set('vert', 'eaVert')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = TEXT_DARK
    run.font.name = "微软雅黑"
    tb.name = "p2-label-text"

print("竖框标签美化完成")

prs.save(SRC)
print("已保存")
