# -*- coding: utf-8 -*-
"""V1.7 第2页恢复到线上(b440648)版本：
1. 删除3个金色竖条(s-bar-1/2/3)
2. 三个标题 13pt加粗 → 12pt 常规
3. 汇总两表 5.71/5.60宽 → 3.42宽原位
4. 竖框标签 2.51高 → 1.47高原位
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

# ===== 1. 删除金色竖条 =====
removed = 0
for sh in list(slide.shapes):
    if sh.name.startswith("s-bar-"):
        slide.shapes._spTree.remove(sh._element)
        removed += 1
print(f"删除金色竖条: {removed} 个")

# ===== 2. 标题恢复 12pt（去掉加粗）=====
titles = [(0.64, 1.32), (0.69, 2.91), (0.69, 4.56)]
for tx, ty in titles:
    for sh in slide.shapes:
        if sh.has_text_frame and abs(sh.left/914400 - tx) < 0.05 and abs(sh.top/914400 - ty) < 0.05:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.bold = False
            print(f"标题({tx},{ty}) 恢复 12pt")

# ===== 3. 汇总两表恢复 =====
tables = [sh for sh in slide.shapes if sh.shape_type == 19]
tbls = [(sh, sh.table) for sh in tables if sh.top/914400 > 4.5]
print(f"找到汇总区表格: {len(tbls)}")

def resize_table(shape, tbl, x_in, w_in, col_widths_in):
    shape.left = Inches(x_in)
    shape.top = Inches(4.94)
    xfrm = shape._element.find(qn('p:xfrm'))
    ext = xfrm.find(qn('a:ext'))
    ext.set('cx', str(int(w_in * 914400)))
    tblGr = tbl._tbl
    grid = tblGr.find(qn('a:tblGrid'))
    gcs = grid.findall(qn('a:gridCol'))
    for gc, w in zip(gcs, col_widths_in):
        gc.set('w', str(int(w * 914400)))
    for tr in tblGr.findall(qn('a:tr')):
        tcs = tr.findall(qn('a:tc'))
        for tc, w in zip(tcs, col_widths_in):
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            tcPr.set('w', str(int(w * 914400)))

# 恢复线上原值
resize_table(*tbls[0], 2.21, 3.42, [1.37, 2.051])
resize_table(*tbls[1], 8.47, 3.42, [1.37, 2.051])
print("两表恢复 3.42in 宽")

# ===== 4. 竖框标签恢复 1.47 高 =====
for sh in slide.shapes:
    if sh.name.startswith("p2-label"):
        x, y = sh.left/914400, sh.top/914400
        if abs(x - 0.90) < 0.05:
            sh.top = Inches(5.22)
            sh.height = Inches(1.47)
        elif abs(x - 7.07) < 0.05:
            sh.top = Inches(5.24)
            sh.height = Inches(1.47)
print("竖框标签恢复 1.47in 高")

prs.save(SRC)
print("已保存")
