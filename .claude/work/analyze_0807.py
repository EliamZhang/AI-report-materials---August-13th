# -*- coding: utf-8 -*-
"""逐页分析 银行流水分享_0807.pptx：所有形状、文字、样式、图表、表格。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn

prs = Presentation("银行流水分享_0807.pptx")
print(f"总页数: {len(prs.slides)}")

for si, slide in enumerate(prs.slides):
    print(f"\n{'='*75}")
    print(f"=== 第{si+1}页 ({len(slide.shapes)}个形状) ===")
    for i, sh in enumerate(slide.shapes):
        x, y = sh.left/914400, sh.top/914400
        w, h = sh.width/914400, sh.height/914400
        st = str(sh.shape_type).split('(')[0].strip()
        line = f"[{i:2d}] {sh.name[:20]:22s}|{st:9s}|({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}"
        # 图表
        if sh.shape_type == 3:
            try:
                chart = sh.chart
                print(f"{line}|CHART type={chart.chart_type} cats={[c for c in chart.plots[0].categories]} series={[(s.name, [v for v in s.values]) for s in chart.plots[0].series]}")
            except Exception as e:
                print(f"{line}|CHART (读取失败 {e})")
            continue
        # 表格
        if sh.shape_type == 19:
            tbl = sh.table
            print(f"{line}|TABLE {len(tbl.rows)}x{len(tbl.columns)}")
            for ri in range(min(len(tbl.rows), 3)):
                cells = [c.text.replace('\n',' ')[:14] for c in tbl.rows[ri].cells]
                print(f"      r{ri}: {cells}")
            continue
        # 文字
        txt = ""
        if sh.has_text_frame:
            parts = []
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else ''
                    b = 'B' if r.font.bold else ''
                    col = ''
                    try:
                        col = f'#{r.font.color.rgb}' if r.font.color and r.font.color.type is not None else ''
                    except: pass
                    parts.append(f"{r.text}({sz}{b}{col})")
            txt = " ".join(parts)
        # 填充
        fill = ""
        try:
            sf = sh._element.spPr.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None:
                    fill = f" FILL:#{clr.get('val')}"
        except: pass
        print(f"{line}{fill} {txt[:80]}")
