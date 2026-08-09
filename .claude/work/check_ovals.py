# -*- coding: utf-8 -*-
"""检查第4页韦恩图圆的XML与实际渲染颜色。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation("银行流水分享_V1.5.pptx")
slide = list(prs.slides)[3]

print("=== 页面形状列表 ===")
for i, sh in enumerate(slide.shapes):
    geom = ""
    try:
        prst = sh._element.spPr.find(qn('a:prstGeom'))
        if prst is not None:
            geom = prst.get('prst')
    except Exception:
        pass
    print(f"[{i}] type={sh.shape_type} name={sh.name!r} pos=({sh.left/914400:.2f},{sh.top/914400:.2f}) size=({sh.width/914400:.2f}x{sh.height/914400:.2f}) geom={geom}")

print()
print("=== 圆的 XML ===")
for sh in slide.shapes:
    if sh.shape_type == 1 and sh.name.startswith("Oval"):
        xml = etree.tostring(sh._element, pretty_print=True).decode()
        # 只打印 spPr 部分
        spPr = sh._element.find(qn('p:spPr'))
        print(etree.tostring(spPr, pretty_print=True).decode())
        print("---")
