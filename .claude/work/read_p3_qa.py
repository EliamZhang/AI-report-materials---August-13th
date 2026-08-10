# -*- coding: utf-8 -*-
"""分析V1.7第3页右侧"AI质检Agent"区域：文字块几何+内容。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[2]

# 右侧区域 x>=6.0, y 1.8-4.4
for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    if x < 6.0:
        continue
    txt = sh.text_frame.text.replace("\n", " ") if sh.has_text_frame else ""
    print(f"[{i:2d}] {sh.name[:22]:26s} ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} 底={y+h:.2f}  {txt[:45]}")
    if sh.has_text_frame and len(txt) > 45:
        print(f"{'':>58s}└ 全文: {txt[:120]}")
