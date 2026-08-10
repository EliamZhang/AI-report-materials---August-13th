# -*- coding: utf-8 -*-
"""V1.7 第3页：压缩文字块字号/段距/行距，使其装进 2.16-4.12 空间。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[2]

target = None
for sh in slide.shapes:
    if abs(sh.left/914400 - 6.90) < 0.05 and abs(sh.top/914400 - 2.16) < 0.05 and sh.has_text_frame:
        target = sh
        break
assert target is not None

tf = target.text_frame
for pi, p in enumerate(tf.paragraphs):
    # 标签行（0,2,4,6）：10.5→10pt，spaceAfter 1→1
    # 正文行（1,3,5,7）：9→8.5pt，spaceAfter 5→3
    is_label = (pi % 2 == 0)
    pPr = p._p.find(qn('a:pPr'))
    if pPr is not None:
        spA = pPr.find(qn('a:spcAft'))
        if spA is not None:
            pts = spA.find(qn('a:spcPts'))
            if pts is not None:
                pts.set('val', '100' if is_label else '300')
    for r in p.runs:
        if is_label:
            r.font.size = Pt(10)
        else:
            r.font.size = Pt(8.5)

print("已压缩：标签10pt，正文8.5pt，段距3pt")

prs.save(SRC)
print("已保存")
