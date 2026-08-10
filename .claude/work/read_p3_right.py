# -*- coding: utf-8 -*-
"""读取V1.7第3页右侧：文字块+两个表格完整内容与样式。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[2]

for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    if x < 5.8:
        continue
    print(f"\n[{i}] {sh.name} ({sh.shape_type}) ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} 底={y+h:.2f}")
    if sh.shape_type == 19:
        tbl = sh.table
        print(f"  {len(tbl.rows)}行x{len(tbl.columns)}列 行高{[round(r.height/914400,2) for r in tbl.rows]} 列宽{[round(c.width/914400,2) for c in tbl.columns]}")
        # 表头样式
        tc = tbl.rows[0].cells[0]
        tcPr = tc._tc.find(qn('a:tcPr'))
        borders = []
        if tcPr is not None:
            for tag, nm in [('a:lnB','底'), ('a:lnT','顶'), ('a:lnL','左'), ('a:lnR','右')]:
                ln = tcPr.find(qn(tag))
                if ln is not None:
                    sf2 = ln.find(qn('a:solidFill'))
                    c = sf2.find(qn('a:srgbClr')).get('val') if sf2 is not None and sf2.find(qn('a:srgbClr')) is not None else None
                    borders.append(f"{nm}{ln.get('w')}/{c}")
            sf = tcPr.find(qn('a:solidFill'))
            fill = sf.find(qn('a:srgbClr')).get('val') if sf is not None and sf.find(qn('a:srgbClr')) is not None else "无"
        print(f"  表头边框[{', '.join(borders)}] 填充={fill}")
        for ri, row in enumerate(tbl.rows):
            cells = []
            for c in row.cells:
                runs = []
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        sz = r.font.size.pt if r.font.size else None
                        b = r.font.bold
                        runs.append(f"{r.text}({sz},{'B' if b else ''})")
                cells.append("".join(runs) if runs else c.text)
            print(f"  r{ri}: {cells}")
    elif sh.has_text_frame:
        for pi, p in enumerate(sh.text_frame.paragraphs):
            runs = []
            for r in p.runs:
                sz = r.font.size.pt if r.font.size else None
                b = r.font.bold
                col = None
                try:
                    col = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
                except Exception:
                    col = "theme"
                runs.append(f"{r.text}({sz}pt,{'B' if b else ''},{col})")
            print(f"  P{pi}: {' '.join(runs)}")
