# -*- coding: utf-8 -*-
"""V1.7 第3页右侧"监控/定位/输出/效果"文字块排版：标签加粗蓝色、正文常规、段距分隔。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[2]

# 形状15 = 文本框 3（监控定位输出效果文字块）
target = None
for sh in slide.shapes:
    if abs(sh.left/914400 - 6.90) < 0.05 and abs(sh.top/914400 - 2.18) < 0.05 and sh.has_text_frame:
        target = sh
        break
assert target is not None, "未找到文字块"

BLUE = RGBColor(0x1F, 0x4E, 0x79)
BODY = RGBColor(0x3D, 0x3C, 0x3A)

# 原内容（逐字保留）
content = [
    ("① 监控", "：捕捉分类批量漂移、同一交易对手归类不一致、覆盖率或命中率明显下降、新描述未被任何引擎识别等异常信号。"),
    ("② 定位", "：对异常逐笔细挖交易描述、金额、方向、交易对手、命中引擎和分类依据，归类为五类原因。"),
    ("③ 输出", "：生成按优先级排序的诊断报告，直接标明问题引擎、需检查的规则、缺失的商户／关键词及应补充的知识库，供业务确认。"),
    ("效果", "：知识库随注册商户的新增、变更、失效持续更新，并从真实交易中补齐注册数据未覆盖的支付描述和商户别名。由此形成“分类 → 监控 → 定位 → 输出 → 更新 → 再分类”的持续优化闭环。"),
]

tf = target.text_frame
tf.clear()
for i, (label, body) in enumerate(content):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    p.space_after = Pt(4) if i < len(content) - 1 else Pt(0)
    r1 = p.add_run()
    r1.text = label
    r1.font.size = Pt(9)
    r1.font.bold = True
    r1.font.color.rgb = BLUE
    r1.font.name = "微软雅黑"
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(9)
    r2.font.bold = False
    r2.font.color.rgb = BODY
    r2.font.name = "微软雅黑"

print("文字块排版完成，4 段标签加粗蓝色")

prs.save(SRC)
print("已保存")
