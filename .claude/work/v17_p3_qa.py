# -*- coding: utf-8 -*-
"""V1.7 第3页"AI质检Agent"文字块：标签独立成行（10.5pt加粗蓝）+ 正文行（9pt灰），纵向4节。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[2]

# 原文字块（形状15：文本框 3）
target = None
for sh in slide.shapes:
    if abs(sh.left/914400 - 6.90) < 0.05 and abs(sh.top/914400 - 2.18) < 0.05 and sh.has_text_frame:
        target = sh
        break
assert target is not None

BLUE = RGBColor(0x1F, 0x4E, 0x79)
BODY = RGBColor(0x3D, 0x3C, 0x3A)

# 原文（逐字保留，标签与正文拆行）
sections = [
    ("① 监控", "捕捉分类批量漂移、同一交易对手归类不一致、覆盖率或命中率明显下降、新描述未被任何引擎识别等异常信号。"),
    ("② 定位", "对异常逐笔细挖交易描述、金额、方向、交易对手、命中引擎和分类依据，归类为五类原因。"),
    ("③ 输出", "生成按优先级排序的诊断报告，直接标明问题引擎、需检查的规则、缺失的商户／关键词及应补充的知识库，供业务确认。"),
    ("效果", "知识库随注册商户的新增、变更、失效持续更新，并从真实交易中补齐注册数据未覆盖的支付描述和商户别名。由此形成“分类 → 监控 → 定位 → 输出 → 更新 → 再分类”的持续优化闭环。"),
]

tf = target.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP

for si, (label, body) in enumerate(sections):
    # 标签行
    p1 = tf.paragraphs[0] if si == 0 else tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    p1.line_spacing = 1.0
    p1.space_before = Pt(0)
    p1.space_after = Pt(1)
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(10.5)
    r1.font.bold = True
    r1.font.color.rgb = BLUE
    r1.font.name = "微软雅黑"
    # 正文行
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.0
    p2.space_before = Pt(0)
    p2.space_after = Pt(5) if si < len(sections) - 1 else Pt(0)
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(9)
    r2.font.bold = False
    r2.font.color.rgb = BODY
    r2.font.name = "微软雅黑"

print("文字块已拆分为 标签行+正文行")

prs.save(SRC)
print("已保存")
