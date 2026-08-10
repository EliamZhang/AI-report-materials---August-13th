# -*- coding: utf-8 -*-
"""V1.7 第2页：交易描述改为更真实的澳洲银行流水风格（两张表同步，金额/日期/方向不变）。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

# 新的交易描述（按日期对应原行，两表共用）
NEW_DESC = [
    "BILL PAY Fair Go Finance A1002345",
    "EFTPOS MCDONALDS BRISBANE QLD",
    "CHEMIST WAREHOUSE 123 BRISBANE QLD",
    "AMZNPRIMEAU*",
    "PAYROLL DEPOSIT ACME PTY LTD",
]

tables = [sh.table for sh in slide.shapes if sh.shape_type == 19]
raw_tbl, proc_tbl = tables[0], tables[1]   # 表1: 原始流水(4列), 表2: 交易明细(9列)

assert len(raw_tbl.rows) == 6 and len(proc_tbl.rows) == 6

for idx, new_desc in enumerate(NEW_DESC):
    row_no = idx + 1
    old_raw = raw_tbl.rows[row_no].cells[3].text
    old_proc = proc_tbl.rows[row_no].cells[1].text
    assert old_raw == old_proc, f"两表描述不一致 r{row_no}: {old_raw} vs {old_proc}"
    raw_tbl.rows[row_no].cells[3].text = new_desc
    proc_tbl.rows[row_no].cells[1].text = new_desc
    print(f"r{row_no}: {old_raw[:30]}... → {new_desc}")

# 分类依据第1行原含“产品关键词”，描述中已无产品词，更新为与实际情况一致的依据
old_basis = proc_tbl.rows[1].cells[8].text
proc_tbl.rows[1].cells[8].text = "交易方向、机构名称及历史交易模式"
print(f"分类依据: {old_basis} → 交易方向、机构名称及历史交易模式")

prs.save(SRC)
print("已保存")
