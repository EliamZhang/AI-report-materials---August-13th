# -*- coding: utf-8 -*-
"""V1.7 第2页：三个小节标题加金色方块徽章（0.30x0.30, #D6A000, 白字数字），标题升13pt加粗。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

GOLD = RGBColor(0xD6, 0xA0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x0A, 0x0A, 0x0A)

# 标题位置与数字
titles = [
    (0.64, 1.32, "1", "1. 原始银行流水（输入示例）"),
    (0.69, 2.91, "2", "2. 交易明细维度：理解每一笔交易（输出示例）"),
    (0.69, 4.56, "3", "3. 用户汇总维度：形成完整负债画像（输出示例）"),
]

SQUARE = 0.30
for tx, ty, num, label in titles:
    # 1. 标题文本框：升 13pt 加粗
    found = None
    for sh in slide.shapes:
        if sh.has_text_frame and abs(sh.left/914400 - tx) < 0.05 and abs(sh.top/914400 - ty) < 0.05:
            found = sh
            break
    assert found is not None, f"未找到 {label}"
    for p in found.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = TEXT_DARK
            r.font.name = "微软雅黑"

    # 2. 金色方块：位于标题左侧，与标题文字垂直居中
    sq_x = tx - SQUARE - 0.05
    sq_y = ty + (found.height/914400 - SQUARE) / 2
    sq = slide.shapes.add_shape(1, Inches(sq_x), Inches(sq_y), Inches(SQUARE), Inches(SQUARE))
    sq.fill.solid()
    sq.fill.fore_color.rgb = GOLD
    sq.line.fill.background()
    sq.shadow.inherit = False
    sq.name = f"s-badge-{num}"

    # 3. 白色数字（覆盖在方块上）
    tb = slide.shapes.add_textbox(Inches(sq_x), Inches(sq_y), Inches(SQUARE), Inches(SQUARE))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "微软雅黑"
    tb.name = f"s-badge-text-{num}"

    print(f"徽章 {num} 方块({sq_x:.2f},{sq_y:.2f}) 标题已升13pt")

prs.save(SRC)
print("已保存")
