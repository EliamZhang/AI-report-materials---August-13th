# -*- coding: utf-8 -*-
"""查看V1.7第2页 用户汇总维度 区域：表格内容、矩形标签与表格的相对位置。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[1]

for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    if sh.shape_type == 19:
        tbl = sh.table
        print(f"[{i}] TABLE pos=({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}in  底={y+h:.2f}  行数={len(tbl.rows)}")
        for ri, row in enumerate(tbl.rows):
            cells = [c.text.replace("\n", " ")[:20] for c in row.cells]
            print(f"    r{ri:2d}: {cells}")
    else:
        txt = sh.text_frame.text.replace("\n", " ")[:30] if sh.has_text_frame else ""
        print(f"[{i}] {sh.shape_type} pos=({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}in 底={y+h:.2f} {txt}")
