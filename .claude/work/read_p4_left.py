# -*- coding: utf-8 -*-
"""读取V1.7第4页左侧区域元素位置（x<5.2in）。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[3]

print("=== 第4页左侧元素 (x < 5.2) ===")
for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    if x < 5.2 or sh.shape_type == 9:
        txt = sh.text_frame.text.replace("\n", " ")[:24] if sh.has_text_frame else ""
        print(f"[{i:2d}] {sh.name[:24]:26s} ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} 底={y+h:.2f}  {txt}")
