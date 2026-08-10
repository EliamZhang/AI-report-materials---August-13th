# -*- coding: utf-8 -*-
"""读取两张用户汇总维度表的列宽、行高与内容（用于重排）。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[1]

for i, sh in enumerate(slide.shapes):
    if sh.shape_type == 19:
        tbl = sh.table
        x, y = sh.left/914400, sh.top/914400
        w = sh.width/914400
        if x > 2.0 and w < 5.0:
            print(f"[{i}] {sh.name} pos=({x:.2f},{y:.2f}) w={w:.2f}in")
            print(f"  列宽: {[round(c.width/914400, 3) for c in tbl.columns]}")
            print(f"  行高: {[round(r.height/914400, 3) for r in tbl.rows]}")
            for ri, row in enumerate(tbl.rows):
                print(f"  r{ri}: {[c.text for c in row.cells]}")
