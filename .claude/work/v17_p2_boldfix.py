# -*- coding: utf-8 -*-
"""修正：三个小节标题恢复加粗（线上为 12pt 加粗）。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[1]

titles = [(0.64, 1.32), (0.69, 2.91), (0.69, 4.56)]
for tx, ty in titles:
    for sh in slide.shapes:
        if sh.has_text_frame and abs(sh.left/914400 - tx) < 0.05 and abs(sh.top/914400 - ty) < 0.05:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
            print(f"标题({tx},{ty}) 恢复加粗")

prs.save(SRC)
print("已保存")
