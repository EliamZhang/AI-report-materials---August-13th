# -*- coding: utf-8 -*-
"""V1.7 第4页左侧垂直拉开：标题上移、圆与圆内数字下移，标签不动。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[3]

def move(sh, dx, dy):
    sh.left = sh.left + Inches(dx)
    sh.top = sh.top + Inches(dy)

# 1. 标题上移 0.52in（TextBox 23 + 蓝条 Rectangle 22）
#    TextBox 23: (0.36,1.67) → (0.36,1.15)
for sh in slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("1. 贷款识别覆盖交集"):
        move(sh, 0, -0.52)
        print(f"标题上移: → ({sh.left/914400:.2f},{sh.top/914400:.2f})")
    elif abs(sh.left/914400 - 0.21) < 0.05 and abs(sh.top/914400 - 1.69) < 0.05 and sh.shape_type == 1:
        move(sh, 0, -0.52)
        print(f"标题蓝条上移: → ({sh.left/914400:.2f},{sh.top/914400:.2f})")

# 2. 两圆下移 0.06in：Oval 63 (0.46,2.09)、Oval 64 (1.76,2.09)
for sh in slide.shapes:
    if sh.shape_type == 1 and sh.name.startswith("Oval") and abs(sh.top/914400 - 2.09) < 0.02:
        move(sh, 0, 0.06)
        print(f"圆下移: {sh.name} → ({sh.left/914400:.2f},{sh.top/914400:.2f})")

# 3. 圆内数字下移 0.06in：TextBox 65(396) 66(79) 67(5284) 68(一致率92.0%)
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t in ("396", "79", "5,284", "一致率 92.0%"):
        move(sh, 0, 0.06)
        print(f"圆内文字下移: {t} → ({sh.left/914400:.2f},{sh.top/914400:.2f})")

prs.save(SRC)
print("已保存")
