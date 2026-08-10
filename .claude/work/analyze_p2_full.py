# -*- coding: utf-8 -*-
"""完整分析V1.7第2页：所有元素位置+垂直关系。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("银行流水分享_V1.7.pptx")
slide = list(prs.slides)[1]

print("=== 第2页 全部元素（按y排序）===")
items = []
for i, sh in enumerate(slide.shapes):
    x, y = sh.left/914400, sh.top/914400
    w, h = sh.width/914400, sh.height/914400
    st = str(sh.shape_type).split('(')[0].strip()
    txt = sh.text_frame.text.replace("\n", " ")[:30] if sh.has_text_frame else ""
    items.append((y, i, sh.name, st, x, y, w, h, txt))

for y, i, name, st, x, yy, w, h, txt in sorted(items):
    print(f"[{i:2d}] {name[:24]:26s} {st:9s} ({x:.2f},{yy:.2f}) {w:.2f}x{h:.2f} 底={yy+h:.2f}  {txt}")
