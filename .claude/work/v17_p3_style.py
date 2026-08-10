# -*- coding: utf-8 -*-
"""V1.7 第3页右侧：两个表格套用V1.3样式（8pt、表头下黑线、数据行灰线、无填充），文字块微调间距。"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[2]

TEXT_DARK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
FONT = "微软雅黑"

def set_cell_borders(tcPr, top=None, bottom=None, left=None, right=None):
    """top/bottom/left/right = (width_emu, color_hex) 或 None"""
    for tag, spec in [('a:lnT', top), ('a:lnB', bottom), ('a:lnL', left), ('a:lnR', right)]:
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        if spec is None:
            continue
        w, color = spec
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set('w', str(w))
        ln.set('cap', 'flat')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr'))
        clr.set('val', color)

def clear_cell_fill(tcPr):
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

def style_table_v13(tbl, font_size=8):
    """套用 V1.3 知识库表格样式"""
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)
    for ri in range(n_rows):
        for ci in range(n_cols):
            cell = tbl.rows[ri].cells[ci]
            tcPr = cell._tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
            clear_cell_fill(tcPr)
            if ri == 0:
                # 表头：底部黑线
                set_cell_borders(tcPr, bottom=(19050, '0A0A0A'))
            else:
                # 数据行：顶部黑线 + 底部细灰线
                set_cell_borders(tcPr, top=(19050, '0A0A0A'), bottom=(9525, 'D6D6D6'))
            # 文字样式
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
                    r.font.color.rgb = TEXT_DARK if ri == 0 else BODY
                    r.font.bold = True if ri == 0 else False
                    # latin + ea
                    rPr = r._r.find(qn('a:rPr'))
                    if rPr is not None:
                        latin = rPr.find(qn('a:latin'))
                        if latin is None:
                            latin = etree.SubElement(rPr, qn('a:latin'))
                        latin.set('typeface', FONT)
                        ea = rPr.find(qn('a:ea'))
                        if ea is None:
                            ea = etree.SubElement(rPr, qn('a:ea'))
                        ea.set('typeface', FONT)
    # 第一列加粗
    for ri in range(1, n_rows):
        cell = tbl.rows[ri].cells[0]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True

# ===== 1. 表格样式 =====
tables = [sh for sh in slide.shapes if sh.shape_type == 19]
right_tables = [sh for sh in tables if sh.left/914400 > 6.0]
print(f"右侧表格数: {len(right_tables)}")
for sh in right_tables:
    style_table_v13(sh.table, font_size=8)
    print(f"已套用样式: {sh.name}")

# ===== 2. 文字块微调间距 =====
tb = None
for sh in slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("① 监控"):
        tb = sh
        break
if tb is not None:
    tf = tb.text_frame
    for pi, p in enumerate(tf.paragraphs):
        p.line_spacing = 1.1
        if pi < len(tf.paragraphs) - 1:
            p.space_after = Pt(3)
    print("文字块间距已微调")

prs.save(SRC)
print("已保存")
