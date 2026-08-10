# -*- coding: utf-8 -*-
"""修复：把统计卡内被误移的 5,284(2.01,5.10) 和 一致率92.0%(2.01,5.44) 移回原位。"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = "银行流水分享_V1.7.pptx"
prs = Presentation(SRC)
slide = list(prs.slides)[3]

fixes = [("5,284", 2.01, 5.10), ("一致率 92.0%", 2.01, 5.44)]
count = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    x, y = sh.left/914400, sh.top/914400
    for ft, fx, fy in fixes:
        if t == ft and abs(x-fx) < 0.05 and abs(y-fy) < 0.05:
            sh.top = sh.top - Inches(0.06)
            count += 1
            print(f"移回: {t} ({x:.2f},{y:.2f}) → ({sh.left/914400:.2f},{sh.top/914400:.2f})")

print(f"共修复 {count} 个")
prs.save(SRC)
print("已保存")
